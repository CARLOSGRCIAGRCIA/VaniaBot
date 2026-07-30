#!/usr/bin/env python3
import sys
import subprocess
import zipfile
import os
from pathlib import Path

MAX_PAGES = 50


def has_extractable_text(doc, sample_pages=3):
    pages_to_check = min(sample_pages, len(doc))
    return any(doc[i].get_text().strip() for i in range(pages_to_check))


def pdf2img(input_path, output_path, fmt='jpeg'):
    import fitz # type: ignore[import]
    doc = fitz.open(input_path)

    if len(doc) > MAX_PAGES:
        doc.close()
        print(f"ERROR: PDF has {len(doc)} pages, max is {MAX_PAGES}", file=sys.stderr)
        sys.exit(3)

    output_dir = Path(output_path).parent
    os.makedirs(output_dir, exist_ok=True)
    base = Path(output_path).stem
    fmt = fmt.lower().replace('jpg', 'jpeg')
    ext = 'png' if fmt == 'png' else 'jpg'

    images = []
    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=200)
        page_path = output_dir / f"{base}_page_{i+1}.{ext}"
        pix.save(str(page_path))
        images.append(str(page_path))

    doc.close()

    if len(images) == 1:
        os.rename(images[0], output_path)
        print(f"OK:{len(images)}:{ext}")
    else:
        zip_path = output_path + '.zip'
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for img_path in images:
                zf.write(img_path, os.path.basename(img_path))
                os.remove(img_path)
        os.rename(zip_path, output_path)
        print(f"OK:{len(images)}:zip")


def pdf2docx(input_path, output_path):
    import fitz # type: ignore[import]
    from pdf2docx import Converter # type: ignore[import]

    doc = fitz.open(input_path)
    page_count = len(doc)

    if page_count > MAX_PAGES:
        doc.close()
        print(f"ERROR: PDF has {page_count} pages, max is {MAX_PAGES}", file=sys.stderr)
        sys.exit(3)

    if not has_extractable_text(doc):
        doc.close()
        print("ERROR: PDF appears to be scanned (no extractable text)", file=sys.stderr)
        sys.exit(2)

    doc.close()

    cv = Converter(input_path)
    try:
        cv.convert(output_path)
    finally:
        cv.close()

    print(f"OK:{page_count}:docx")


def pdf2ppt(input_path, output_path):
    import fitz # type: ignore[import]
    from pptx import Presentation # type: ignore[import]
    from pptx.util import Emu # type: ignore[import]

    doc = fitz.open(input_path)

    if len(doc) == 0:
        doc.close()
        print("ERROR: PDF has no pages", file=sys.stderr)
        sys.exit(1)

    if len(doc) > MAX_PAGES:
        doc.close()
        print(f"ERROR: PDF has {len(doc)} pages, max is {MAX_PAGES}", file=sys.stderr)
        sys.exit(3)

    prs = Presentation()
    first_rect = doc[0].rect
    prs.slide_width = Emu(int(first_rect.width * 12700))
    prs.slide_height = Emu(int(first_rect.height * 12700))
    blank_layout = prs.slide_layouts[6]

    tmp_dir = Path(output_path).parent / f"_pdf2ppt_{os.getpid()}"
    os.makedirs(tmp_dir, exist_ok=True)

    page_count = len(doc)

    try:
        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=150)
            img_path = tmp_dir / f"page_{i+1}.png"
            pix.save(str(img_path))

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(img_path), 0, 0,
                width=prs.slide_width, height=prs.slide_height
            )
            os.remove(img_path)

        prs.save(output_path)
    finally:
        doc.close()
        try:
            os.rmdir(tmp_dir)
        except OSError:
            pass

    print(f"OK:{page_count}:pptx")


def libreoffice_convert(input_path, output_path):
    output_dir = str(Path(output_path).parent)
    result = subprocess.run([
        'libreoffice', '--headless', '--convert-to', 'pdf',
        '--outdir', output_dir, input_path
    ], capture_output=True, text=True, timeout=120)

    if result.returncode != 0:
        print(f"ERROR: {result.stderr.strip()}", file=sys.stderr)
        sys.exit(1)

    input_name = Path(input_path).stem
    expected = Path(output_dir) / f"{input_name}.pdf"
    if expected.exists():
        os.rename(str(expected), output_path)

    print("OK:1:pdf")


if __name__ == '__main__':
    if len(sys.argv) < 4:
        print("ERROR: Usage: bridge.py <action> <input> <output> [format]", file=sys.stderr)
        sys.exit(1)

    action = sys.argv[1]
    input_path = sys.argv[2]
    output_path = sys.argv[3]
    fmt = sys.argv[4] if len(sys.argv) > 4 else 'jpeg'

    try:
        if action == 'pdf2img':
            pdf2img(input_path, output_path, fmt)
        elif action == 'ppt2pdf':
            libreoffice_convert(input_path, output_path)
        elif action == 'docx2pdf':
            libreoffice_convert(input_path, output_path)
        elif action == 'pdf2docx':
            pdf2docx(input_path, output_path)
        elif action == 'pdf2ppt':
            pdf2ppt(input_path, output_path)
        else:
            print(f"ERROR: Unknown action: {action}", file=sys.stderr)
            sys.exit(1)
    except SystemExit:
        raise
    except Exception as e:
        print(f"ERROR: {e}", file=sys.stderr)
        sys.exit(1)